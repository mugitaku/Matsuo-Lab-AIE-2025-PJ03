import os
import tempfile
from typing import List, Dict, Any, Optional
from pptx import Presentation
import PyPDF2
from pdf2image import convert_from_path
import base64
from io import BytesIO
from PIL import Image


class SlideContent:
    def __init__(self, slide_number: int, text_content: str, image_content: Optional[bytes] = None):
        self.slide_number = slide_number
        self.text_content = text_content
        self.image_content = image_content
        self.image_base64 = None
        if image_content:
            self.image_base64 = base64.b64encode(image_content).decode('utf-8')


class FileParser:
    def __init__(self):
        self.supported_formats = ['.pptx', '.ppt', '.pdf']
    
    def parse_file(self, file_path: str) -> List[SlideContent]:
        file_ext = os.path.splitext(file_path)[1].lower()
        
        if file_ext not in self.supported_formats:
            raise ValueError(f"Unsupported file format: {file_ext}")
        
        if file_ext in ['.pptx', '.ppt']:
            return self._parse_powerpoint(file_path)
        elif file_ext == '.pdf':
            return self._parse_pdf(file_path)
    
    def _parse_powerpoint(self, file_path: str) -> List[SlideContent]:
        presentation = Presentation(file_path)
        slides_content = []
        
        for idx, slide in enumerate(presentation.slides, 1):
            text_content = self._extract_text_from_slide(slide)
            image_content = self._extract_image_from_slide(slide, idx)
            slides_content.append(SlideContent(idx, text_content, image_content))
        
        return slides_content
    
    def _extract_text_from_slide(self, slide) -> str:
        text_parts = []
        
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_parts.append(shape.text)
            
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        text_parts.append(cell.text)
        
        return "\n".join(text_parts)
    
    def _extract_image_from_slide(self, slide, slide_number: int) -> Optional[bytes]:
        try:
            img_buffer = BytesIO()
            
            # Create a simple representation of the slide
            # In a real implementation, we would render the slide properly
            img = Image.new('RGB', (1024, 768), color='white')
            img.save(img_buffer, format='PNG')
            
            return img_buffer.getvalue()
        except Exception:
            return None
    
    def _parse_pdf(self, file_path: str) -> List[SlideContent]:
        slides_content = []
        
        # Extract text from PDF
        pdf_reader = PyPDF2.PdfReader(file_path)
        
        # Convert PDF pages to images
        try:
            images = convert_from_path(file_path, dpi=150)
        except Exception:
            images = []
        
        for idx, page in enumerate(pdf_reader.pages):
            text_content = page.extract_text()
            
            # Get corresponding image if available
            image_content = None
            if idx < len(images):
                img_buffer = BytesIO()
                images[idx].save(img_buffer, format='PNG')
                image_content = img_buffer.getvalue()
            
            slides_content.append(SlideContent(idx + 1, text_content, image_content))
        
        return slides_content
    
    def extract_metadata(self, file_path: str) -> Dict[str, Any]:
        metadata = {
            'file_name': os.path.basename(file_path),
            'file_size': os.path.getsize(file_path),
            'file_type': os.path.splitext(file_path)[1].lower()
        }
        
        file_ext = metadata['file_type']
        
        if file_ext in ['.pptx', '.ppt']:
            try:
                presentation = Presentation(file_path)
                metadata['slide_count'] = len(presentation.slides)
                if hasattr(presentation.core_properties, 'title'):
                    metadata['title'] = presentation.core_properties.title
                if hasattr(presentation.core_properties, 'author'):
                    metadata['author'] = presentation.core_properties.author
            except Exception:
                pass
        
        elif file_ext == '.pdf':
            try:
                pdf_reader = PyPDF2.PdfReader(file_path)
                metadata['page_count'] = len(pdf_reader.pages)
                if pdf_reader.metadata:
                    metadata['title'] = pdf_reader.metadata.get('/Title', '')
                    metadata['author'] = pdf_reader.metadata.get('/Author', '')
            except Exception:
                pass
        
        return metadata